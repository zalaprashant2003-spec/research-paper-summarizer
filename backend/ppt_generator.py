from pptx import Presentation
from pptx.util import Inches, Pt
import uuid


def create_slide(prs, title, bullets):

    slide = prs.slides.add_slide(
        prs.slide_layouts[5]
    )

    slide.shapes.title.text = title

    textbox = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(1.2),
        Inches(8.5),
        Inches(4)
    )

    frame = textbox.text_frame

    frame.clear()

    for bullet in bullets:

        p = frame.add_paragraph()

        p.text = f"• {bullet}"

        p.font.size = Pt(20)

    return slide


def create_ppt(data):

    prs = Presentation()

    title_slide = prs.slides.add_slide(
        prs.slide_layouts[0]
    )

    title_slide.shapes.title.text = (
        "Research Paper Analysis"
    )

    slides_data = [
        ("Problem", data["problem"]),
        ("Method", data["method"]),
        ("Results", data["results"]),
        ("Conclusion", data["conclusion"]),
        ("Future Work", data["future_work"]),
    ]

    for title, bullets in slides_data:

        create_slide(
            prs,
            title,
            bullets
        )

    file_name = (
        f"{uuid.uuid4()}.pptx"
    )

    prs.save(file_name)

    return file_name